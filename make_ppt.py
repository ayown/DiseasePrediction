"""
Disease Prediction System — Premium PPT v3
Fixes: all text-overflow / overlap issues, refined dark UI, working animations.
"""

from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ── Output path ───────────────────────────────────────────────────────────────
BASE = Path(__file__).resolve().parent
OUT  = BASE / "Disease_Prediction_PPT.pptx"

# ── Slide canvas ──────────────────────────────────────────────────────────────
SW, SH = 13.33, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ── Colour palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x07, 0x0B, 0x14)   # slide background
PANEL    = RGBColor(0x10, 0x15, 0x22)   # card fill
PANEL2   = RGBColor(0x17, 0x1D, 0x2C)   # lighter card
BORDER   = RGBColor(0x28, 0x2F, 0x42)   # card stroke
WHITE    = RGBColor(0xF2, 0xF5, 0xFA)
MUTED    = RGBColor(0x8A, 0x94, 0xA6)
SOFT     = RGBColor(0xC8, 0xD0, 0xDE)
ORANGE   = RGBColor(0xFF, 0x6B, 0x2B)
AMBER    = RGBColor(0xFF, 0xB3, 0x30)
TEAL     = RGBColor(0x2D, 0xCE, 0xBF)
BLUE     = RGBColor(0x5B, 0xA8, 0xFF)
GREEN    = RGBColor(0x3D, 0xD6, 0x8C)
RED      = RGBColor(0xF9, 0x6B, 0x6B)
PURPLE   = RGBColor(0xBD, 0x83, 0xFF)
LIME     = RGBColor(0xA8, 0xE6, 0x3D)

# ── Layout constants ──────────────────────────────────────────────────────────
ML = 0.5    # left margin
MR = 0.5    # right margin
MT = 1.22   # top of content (below header)
MB = 0.42   # footer height
CW = SW - ML - MR          # 12.33"  content width
CH = SH - MT - MB - 0.10   # ≈5.78"  content height

# ═══════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL DRAW HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def _shape(slide, mso, x, y, w, h):
    return slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))


def box(slide, x, y, w, h, fill=PANEL, stroke=None, sw=1.0, rounded=True, shadow=False):
    """Rectangle / rounded-rectangle with optional shadow."""
    mso = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = _shape(slide, mso, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if stroke:
        s.line.color.rgb = stroke
        s.line.width = Pt(sw)
    else:
        s.line.fill.background()
    if shadow:
        _drop_shadow(s)
    return s


def oval(slide, x, y, w, h, fill=ORANGE, stroke=None):
    s = _shape(slide, MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if stroke:
        s.line.color.rgb = stroke
    else:
        s.line.fill.background()
    return s


def connector(slide, x1, y1, x2, y2, color=BORDER, width=1.0):
    s = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = color
    s.line.width = Pt(width)
    return s


def _drop_shadow(shape):
    sp_pr = shape._element.find(qn('p:spPr'))
    if sp_pr is None:
        return
    eff = sp_pr.find(qn('a:effectLst'))
    if eff is None:
        eff = etree.SubElement(sp_pr, qn('a:effectLst'))
    shdw = etree.SubElement(eff, qn('a:outerShdw'))
    shdw.set('blurRad', '40000')
    shdw.set('dist', '12000')
    shdw.set('dir', '5400000')
    shdw.set('algn', 'tl')
    shdw.set('rotWithShape', '0')
    c = etree.SubElement(shdw, qn('a:srgbClr'))
    c.set('val', '000000')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', '55000')


def txt(slide, value, x, y, w, h,
        size=14, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font='Segoe UI', valign=MSO_ANCHOR.MIDDLE):
    """
    Add a text box. h is the RESERVED height — always make it larger than needed.
    Never use AUTO_SIZE; instead size explicitly to prevent overlap.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap    = True
    tf.vertical_anchor = valign
    # NO auto-size — prevents runaway text
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = value
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    run.font.name       = font
    return tb


def txt_multi(slide, lines, x, y, w, h,
              size=13, color=WHITE, bold=False,
              align=PP_ALIGN.LEFT, font='Segoe UI',
              line_spacing=1.15, para_space=4):
    """
    Multi-line text box with explicit per-line paragraphs.
    lines: list of str
    """
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt as _Pt
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap     = True
    tf.auto_size     = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for li, line_text in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment   = align
        p.space_after = _Pt(para_space)
        run = p.add_run()
        run.text           = line_text
        run.font.size      = _Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = font
    return tb


# ═══════════════════════════════════════════════════════════════════════════════
# REUSABLE COMPONENTS
# ═══════════════════════════════════════════════════════════════════════════════

def slide_bg(slide):
    """Fill slide background with BG colour."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG


def top_strip(slide, color=ORANGE, height=0.055):
    box(slide, 0, 0, SW, height, fill=color, rounded=False)


def header_block(slide, eyebrow, title, subtitle=None):
    """Standard header: coloured pill + big title + optional subtitle."""
    # eyebrow pill
    box(slide, ML, 0.14, len(eyebrow) * 0.115 + 0.40, 0.30,
        fill=ORANGE, rounded=True)
    txt(slide, eyebrow.upper(), ML + 0.14, 0.15, 3.5, 0.26,
        size=8.5, color=BG, bold=True, font='Segoe UI Semibold',
        align=PP_ALIGN.LEFT)
    # title
    txt(slide, title, ML, 0.50, CW, 0.55,
        size=28, color=WHITE, bold=True, font='Segoe UI Semibold')
    # subtitle
    if subtitle:
        txt(slide, subtitle, ML, 0.98, CW * 0.82, 0.22,
            size=10, color=MUTED, font='Segoe UI')


def footer(slide):
    box(slide, 0, SH - MB + 0.06, SW, 0.01, fill=BORDER, rounded=False)
    txt(slide, 'Disease Prediction System  ·  Using Machine Learning',
        ML, SH - MB + 0.10, 6.0, 0.25, size=7.5, color=MUTED)
    txt(slide, 'MLC-2  |  CSE 3968  |  ITER, SOA University  |  May 2026',
        SW - 5.5, SH - MB + 0.10, 5.0, 0.25,
        size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, accent=ORANGE, fill=PANEL, stroke=BORDER, shadow=True):
    """Glass card: shadow layer + filled card + left accent bar."""
    # shadow
    if shadow:
        box(slide, x + 0.06, y + 0.07, w, h, fill=RGBColor(0x02, 0x04, 0x09),
            rounded=True, stroke=None)
    # body
    c = box(slide, x, y, w, h, fill=fill, stroke=stroke, rounded=True)
    # left accent strip
    box(slide, x, y + 0.06, 0.055, h - 0.12, fill=accent, rounded=False)
    return c


def metric_tile(slide, x, y, w, h, value, label, accent=ORANGE):
    """KPI tile: big coloured number + muted label."""
    card(slide, x, y, w, h, accent=accent, fill=PANEL2)
    # value occupies top 44% of tile; label starts at 62% — clear gap between them
    txt(slide, value, x + 0.14, y + 0.08, w - 0.26, h * 0.44,
        size=26, color=accent, bold=True, font='Segoe UI Black',
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)
    txt(slide, label, x + 0.10, y + h * 0.62, w - 0.20, h * 0.34,
        size=9, color=MUTED, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.TOP)


def progress_row(slide, label, pct, x, y, w, accent=ORANGE, max_bar=8.5):
    """Label + dark track + filled bar + value."""
    lw = 2.6   # label column width
    bw = max_bar  # bar track width
    vw = 0.85  # value column width
    # label
    txt(slide, label, x, y + 0.02, lw, 0.26, size=10.5, color=SOFT,
        bold=False, font='Segoe UI')
    # track
    box(slide, x + lw + 0.12, y + 0.07, bw, 0.14,
        fill=RGBColor(0x1E, 0x26, 0x38), rounded=True, stroke=None)
    # fill
    fill_w = bw * pct / 100
    box(slide, x + lw + 0.12, y + 0.07, fill_w, 0.14,
        fill=accent, rounded=True, stroke=None, shadow=True)
    # value
    txt(slide, f'{pct:.2f}%', x + lw + bw + 0.22, y + 0.01, vw, 0.26,
        size=10.5, color=accent, bold=True, align=PP_ALIGN.RIGHT)


def tag_pill(slide, label, x, y, accent=ORANGE):
    pw = len(label) * 0.095 + 0.36
    box(slide, x, y, pw, 0.27, fill=accent, rounded=True, stroke=None, shadow=True)
    txt(slide, label, x + 0.12, y + 0.04, pw - 0.18, 0.18,
        size=7.8, color=BG, bold=True, align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
    return pw


def num_badge(slide, num, x, y, accent=ORANGE, d=0.42):
    oval(slide, x, y, d, d, fill=accent)
    txt(slide, str(num), x, y + d * 0.2, d, d * 0.55,
        size=11, color=BG, bold=True, align=PP_ALIGN.CENTER, font='Segoe UI Black')


def _set_rounded_corners(shape, emu=50000):
    """Make rounded rectangle have tighter radius."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    else:
        for gd in list(avLst):
            avLst.remove(gd)
    gd = etree.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {emu}')


# ═══════════════════════════════════════════════════════════════════════════════
# ANIMATION ENGINE  (auto-play, sequential fade-in)
# ═══════════════════════════════════════════════════════════════════════════════

class Animator:
    """
    Builds a p:timing tree on a slide with sequential afterPrev fade animations.
    Call .add(shape, delay_ms) for each element, then .apply().
    """
    def __init__(self, slide):
        self.slide   = slide
        self.entries = []   # (shape_id, delay_ms, dur_ms, effect)
        self._id     = 100

    def _nid(self):
        self._id += 1
        return str(self._id)

    def add(self, shape, delay_ms=0, dur_ms=500, effect='fade'):
        self.entries.append((shape.shape_id, delay_ms, dur_ms, effect))
        return self

    def apply(self):
        if not self.entries:
            return
        sld = self.slide._element
        # Remove existing timing
        for el in sld.findall(qn('p:timing')):
            sld.remove(el)

        timing     = etree.SubElement(sld, qn('p:timing'))
        tnLst      = etree.SubElement(timing, qn('p:tnLst'))
        par_root   = etree.SubElement(tnLst, qn('p:par'))
        cTn_root   = etree.SubElement(par_root, qn('p:cTn'))
        cTn_root.set('id', self._nid())
        cTn_root.set('dur', 'indefinite')
        cTn_root.set('restart', 'whenNotActive')
        cTn_root.set('nodeType', 'tmRoot')
        chLst_root = etree.SubElement(cTn_root, qn('p:childTnLst'))

        seq        = etree.SubElement(chLst_root, qn('p:seq'))
        seq.set('concurrent', '1')
        seq.set('nextAc', 'seek')
        cTn_seq    = etree.SubElement(seq, qn('p:cTn'))
        seq_id     = self._nid()
        cTn_seq.set('id', seq_id)
        cTn_seq.set('dur', 'indefinite')
        cTn_seq.set('nodeType', 'mainSeq')
        chLst_seq  = etree.SubElement(cTn_seq, qn('p:childTnLst'))

        # nextCondLst for the seq
        nc     = etree.SubElement(seq, qn('p:nextCondLst'))
        cond_n = etree.SubElement(nc, qn('p:cond'))
        cond_n.set('evt', 'onClick')
        cond_n.set('delay', '0')
        tn_ref = etree.SubElement(cond_n, qn('p:tn'))
        tn_ref.set('val', seq_id)

        FILTER_MAP = {'fade': 'fade', 'appear': 'appear',
                      'fly':  'flyInFromBottom', 'wipe': 'wipe'}
        PRESET_MAP = {'fade': 10, 'appear': 1, 'fly': 2, 'wipe': 56}

        for spid, delay, dur, effect in self.entries:
            filt   = FILTER_MAP.get(effect, 'fade')
            preset = PRESET_MAP.get(effect, 10)

            # outer par (afterEffect = auto after prev)
            p_outer  = etree.SubElement(chLst_seq, qn('p:par'))
            cTn_out  = etree.SubElement(p_outer, qn('p:cTn'))
            cTn_out.set('id', self._nid())
            cTn_out.set('fill', 'hold')
            cTn_out.set('nodeType', 'afterEffect')
            stCL_out = etree.SubElement(cTn_out, qn('p:stCondLst'))
            cond_out = etree.SubElement(stCL_out, qn('p:cond'))
            cond_out.set('delay', str(delay))
            chLst_out = etree.SubElement(cTn_out, qn('p:childTnLst'))

            # inner par
            p_in    = etree.SubElement(chLst_out, qn('p:par'))
            cTn_in  = etree.SubElement(p_in, qn('p:cTn'))
            cTn_in.set('id', self._nid())
            cTn_in.set('fill', 'hold')
            stCL_in = etree.SubElement(cTn_in, qn('p:stCondLst'))
            cond_in = etree.SubElement(stCL_in, qn('p:cond'))
            cond_in.set('delay', '0')
            chLst_in = etree.SubElement(cTn_in, qn('p:childTnLst'))

            # effect par
            p_ef     = etree.SubElement(chLst_in, qn('p:par'))
            cTn_ef   = etree.SubElement(p_ef, qn('p:cTn'))
            cTn_ef.set('id', self._nid())
            cTn_ef.set('presetID', str(preset))
            cTn_ef.set('presetClass', 'entr')
            cTn_ef.set('presetSubtype', '0')
            cTn_ef.set('fill', 'hold')
            cTn_ef.set('grpId', '0')
            cTn_ef.set('nodeType', 'withEffect')
            stCL_ef  = etree.SubElement(cTn_ef, qn('p:stCondLst'))
            cond_ef  = etree.SubElement(stCL_ef, qn('p:cond'))
            cond_ef.set('delay', '0')
            chLst_ef = etree.SubElement(cTn_ef, qn('p:childTnLst'))

            # animEffect
            aef     = etree.SubElement(chLst_ef, qn('p:animEffect'))
            aef.set('transition', 'in')
            aef.set('filter', filt)
            cBhvr   = etree.SubElement(aef, qn('p:cBhvr'))
            cTn_bh  = etree.SubElement(cBhvr, qn('p:cTn'))
            cTn_bh.set('id', self._nid())
            cTn_bh.set('dur', str(dur))
            stCL_bh = etree.SubElement(cTn_bh, qn('p:stCondLst'))
            cond_bh = etree.SubElement(stCL_bh, qn('p:cond'))
            cond_bh.set('delay', '0')
            tgtEl   = etree.SubElement(cBhvr, qn('p:tgtEl'))
            spTgt   = etree.SubElement(tgtEl, qn('p:spTgt'))
            spTgt.set('spid', str(spid))


def transition(slide, kind='fade', dur=650):
    sld = slide._element
    for el in sld.findall(qn('p:transition')):
        sld.remove(el)
    tr = etree.SubElement(sld, qn('p:transition'))
    tr.set('dur', str(dur))
    if kind == 'push':
        e = etree.SubElement(tr, qn('p:push'))
        e.set('dir', 'l')
    elif kind == 'wipe':
        e = etree.SubElement(tr, qn('p:wipe'))
        e.set('dir', 'l')
    elif kind == 'zoom':
        e = etree.SubElement(tr, qn('p:zoom'))
        e.set('dir', 'in')
    else:
        e = etree.SubElement(tr, qn('p:fade'))
        e.set('thruBlk', '0')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── 1. COVER ──────────────────────────────────────────────────────────────────
def s_cover():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    anim = Animator(slide)

    # Decorative glows (large translucent circles)
    g1 = oval(slide, 7.8, -1.2, 5.2, 5.2, fill=RGBColor(0x40, 0x15, 0x05))
    g2 = oval(slide, 9.6,  3.8, 2.8, 2.8, fill=RGBColor(0x35, 0x2A, 0x04))
    g3 = oval(slide, -0.9, 4.5, 2.4, 2.4, fill=RGBColor(0x08, 0x26, 0x36))

    # Top colour bar
    top = box(slide, 0, 0, SW, 0.055, fill=ORANGE, rounded=False)

    # Left content ────────────────────────────────────────────────────
    pill_w = tag_pill(slide, 'AI HEALTHCARE  ·  ML PROJECT', 0.55, 0.76)

    # Each block sized conservatively so bottom edge never reaches next block's top
    h1 = txt(slide, 'Disease',    0.55, 1.30, 6.4, 0.78,
             size=52, bold=True, font='Segoe UI Black')
    # h1 textbox bottom = 2.08
    h2 = txt(slide, 'Prediction', 0.55, 2.16, 7.2, 0.92,
             size=60, color=ORANGE, bold=True, font='Segoe UI Black')
    # h2 textbox bottom = 3.08
    h3 = txt(slide, 'System',     0.55, 3.16, 6.0, 0.68,
             size=40, color=SOFT, bold=True, font='Segoe UI Semibold')
    # h3 textbox bottom = 3.84

    connector(slide, 0.55, 3.96, 5.2, 3.96, ORANGE, 1.8)

    sub = txt(slide,
              'A severity-aware, age-context disease classification pipeline '
              'comparing Decision Tree, Random Forest, XGBoost, and Stacking Ensemble '
              'across 41 disease categories.',
              0.58, 4.10, 5.65, 0.78,
              size=13, color=MUTED)
    # sub bottom = 4.88

    # KPI row — starts at 5.05 (sub bottom ≈4.88, gap 0.17")
    kpis = [('93.50%', 'Best Accuracy', ORANGE),
            ('41',     'Diseases',      AMBER),
            ('135',    'Features',      TEAL)]
    for i, (v, l, c) in enumerate(kpis):
        metric_tile(slide, 0.55 + i * 1.98, 5.05, 1.82, 1.08, v, l, c)

    # Right panel ─────────────────────────────────────────────────────
    rpanel = card(slide, 7.90, 0.98, 4.88, 5.58, accent=ORANGE, fill=PANEL2, shadow=True)

    txt(slide, 'ML Pipeline', 8.28, 1.34, 3.6, 0.35,
        size=15, bold=True, font='Segoe UI Semibold')

    steps = [('01', 'Symptoms',          ORANGE),
             ('02', 'Severity Encoding',  AMBER),
             ('03', 'Age Features',       TEAL),
             ('04', 'Ensemble Training',  BLUE),
             ('05', 'Prediction Output',  GREEN)]
    for i, (n, lbl, c) in enumerate(steps):
        y = 1.86 + i * 0.78
        oval(slide, 8.30, y, 0.38, 0.38, fill=c)
        txt(slide, n,   8.30, y + 0.07, 0.38, 0.22,
            size=8, color=BG, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, lbl, 8.82, y + 0.06, 3.5, 0.28,
            size=11.5, color=SOFT, bold=True, font='Segoe UI Semibold')
        if i < 4:
            connector(slide, 8.49, y + 0.40, 8.49, y + 0.74, BORDER, 1.1)

    txt(slide, "ITER, SOA University  |  May 2026",
        8.28, 6.38, 4.18, 0.24, size=9, color=MUTED)

    footer(slide)
    anim.add(h1, 0, 550, 'fade').add(h2, 120, 550, 'fade') \
        .add(h3, 240, 500, 'fade').add(sub, 400, 500, 'fade') \
        .add(rpanel, 200, 600, 'fade').apply()
    transition(slide, 'fade', 900)


# ── 2. AGENDA ─────────────────────────────────────────────────────────────────
def s_agenda():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, AMBER)
    header_block(slide, 'Presentation Flow', 'Agenda', 'From data to validated predictions — a complete walkthrough.')
    footer(slide)
    anim = Animator(slide)

    items = [
        ('01', 'Problem Statement', 'Why ML for disease prediction?',      ORANGE),
        ('02', 'Dataset & Features','4,925 records · 41 diseases · 135 features', TEAL),
        ('03', 'Pipeline',          'End-to-end preprocessing flow',       BLUE),
        ('04', 'Four Models',       'DT · RF · XGB · Stacking Ensemble',   AMBER),
        ('05', 'Results',           'Accuracy, precision, recall, F1',     GREEN),
        ('06', 'Insights & Future', 'Feature importance · next steps',     PURPLE),
    ]

    # 2 × 3 grid
    cols, col_x = 3, [ML, ML + 4.18, ML + 8.36]
    for i, (num, title, sub, accent) in enumerate(items):
        col, row = i % cols, i // cols
        x = col_x[col]
        y = MT + row * 2.42

        c = card(slide, x, y, 3.78, 2.12, accent=accent, fill=PANEL, shadow=True)
        num_badge(slide, num, x + 0.22, y + 0.28, accent)
        txt(slide, title, x + 0.82, y + 0.26, 2.6, 0.36,
            size=13.5, bold=True, color=WHITE, font='Segoe UI Semibold')
        txt(slide, sub,   x + 0.22, y + 0.80, 3.3, 0.52,
            size=10.5, color=MUTED)
        anim.add(c, i * 110, 480, 'fade')

    anim.apply()
    transition(slide, 'push', 650)


# ── 3. PROBLEM STATEMENT ──────────────────────────────────────────────────────
def s_problem():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, ORANGE)
    header_block(slide, 'Problem', 'Clinical Decision Support Needs Better Signal',
                 'Reframing symptom-based diagnosis as a multi-class classification problem.')
    footer(slide)
    anim = Animator(slide)

    # Quote box (full width)
    qb = box(slide, ML, MT, CW, 0.82,
             fill=RGBColor(0x1A, 0x10, 0x06), stroke=ORANGE, sw=1.4, rounded=True)
    txt(slide,
        '"A patient reports: fever · joint pain · fatigue. '
        'Is it Malaria, Dengue, Typhoid, or Arthritis? '
        'The same symptoms → 4+ possible diagnoses."',
        ML + 0.28, MT + 0.12, CW - 0.56, 0.56,
        size=14.5, color=WHITE, italic=True, align=PP_ALIGN.CENTER,
        font='Segoe UI Italic')
    anim.add(qb, 0, 500)

    # 3 pillars
    pillars = [
        ('Challenge',
         ['• 41 overlapping disease classes',
          '• Symptoms shared across diseases',
          '• Hard to distinguish without lab tests'],
         RED),
        ('Feature Gap',
         ['• Binary encoding loses severity nuance',
          '• No age-context in prior systems',
          '• Noise from incomplete patient reports'],
         AMBER),
        ('Our Solution',
         ['• Severity-weighted feature matrix',
          '• Epidemiological age feature added',
          '• Stacking Ensemble meta-learner'],
         GREEN),
    ]

    pw = (CW - 0.50) / 3
    for i, (heading, bullets, accent) in enumerate(pillars):
        x = ML + i * (pw + 0.25)
        y = MT + 1.02

        c = card(slide, x, y, pw, 3.38, accent=accent, fill=PANEL, shadow=True)
        # heading strip
        box(slide, x, y, pw, 0.42, fill=accent, rounded=True, stroke=None)
        txt(slide, heading, x, y + 0.02, pw, 0.36,
            size=13, bold=True, color=BG,
            align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
        # bullets
        txt_multi(slide, bullets, x + 0.20, y + 0.55, pw - 0.28, 2.65,
                  size=11.5, color=SOFT, line_spacing=1.1, para_space=8)
        anim.add(c, 120 + i * 150, 500, 'fly')

    # Bottom task strip
    tb = box(slide, ML, SH - MB - 0.82, CW, 0.54,
             fill=RGBColor(0x12, 0x19, 0x2A), stroke=TEAL, sw=1.2, rounded=True)
    txt(slide,
        'Task: predict which of 41 diseases a patient most likely has '
        'from a 135-feature vector (severity-weighted symptoms + age)',
        ML + 0.25, SH - MB - 0.76, CW - 0.50, 0.42,
        size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    anim.add(tb, 600, 480)

    anim.apply()
    transition(slide, 'fade', 700)


# ── 4. DATASET & FEATURES ─────────────────────────────────────────────────────
def s_dataset():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, TEAL)
    header_block(slide, 'Dataset', 'Feature Engineering Layer',
                 'Severity-weighted symptoms + epidemiological age context → 135 model-ready features.')
    footer(slide)
    anim = Animator(slide)

    # KPI row (5 tiles)
    kpis = [('4,925', 'Records',        ORANGE),
            ('41',    'Disease Classes', TEAL),
            ('136',   'Symptom Cols',    AMBER),
            ('+1',    'Age Feature',     BLUE),
            ('135',   'Final Features',  GREEN)]
    tw = (CW - 0.40) / 5
    for i, (v, l, c) in enumerate(kpis):
        t = box(slide, ML + i * (tw + 0.10), MT, tw, 0.95,
                fill=PANEL2, stroke=BORDER, rounded=True)
        txt(slide, v, ML + i * (tw + 0.10), MT + 0.06,
            tw, 0.45, size=22, color=c, bold=True,
            align=PP_ALIGN.CENTER, font='Segoe UI Black')
        txt(slide, l, ML + i * (tw + 0.10), MT + 0.54,
            tw, 0.32, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        anim.add(t, i * 90, 400, 'fade')

    # Left panel — encoding
    lc = card(slide, ML, MT + 1.12, 5.78, 3.40, accent=ORANGE, fill=PANEL, shadow=True)
    txt(slide, 'Severity-Weighted Encoding', ML + 0.22, MT + 1.30, 5.3, 0.34,
        size=16, bold=True, color=WHITE, font='Segoe UI Semibold')
    enc_rows = [
        ('Raw Input',  'Up to 17 symptom names per patient record'),
        ('Lookup',     'Severity dictionary: each symptom → weight 1–7'),
        ('Encode',     'Build 136-col matrix; 0 = symptom absent'),
        ('Dropout',    'Randomly zero 50% of non-zero entries (seed 42)'),
        ('Filter',     'Remove zero-variance cols → 135 final features'),
    ]
    for i, (a, b) in enumerate(enc_rows):
        y = MT + 1.78 + i * 0.52
        box(slide, ML + 0.20, y, 5.3, 0.42,
            fill=RGBColor(0x16, 0x1E, 0x30), stroke=None, rounded=True)
        txt(slide, a, ML + 0.32, y + 0.09, 1.0, 0.25,
            size=8.5, color=ORANGE, bold=True, font='Segoe UI Semibold')
        txt(slide, b, ML + 1.42, y + 0.09, 3.95, 0.25,
            size=9.5, color=SOFT)
    anim.add(lc, 150, 500, 'fade')

    # Right panel — age feature
    rc = card(slide, ML + 6.12, MT + 1.12, 6.12, 3.40, accent=TEAL, fill=PANEL, shadow=True)
    txt(slide, 'Age-Aware Feature  (Novel Addition)',
        ML + 6.34, MT + 1.30, 5.65, 0.34,
        size=16, bold=True, color=WHITE, font='Segoe UI Semibold')
    txt(slide, 'Each disease has an epidemiological age range. '
        'A patient age is sampled from that range — giving the model demographic prior knowledge.',
        ML + 6.34, MT + 1.74, 5.65, 0.62,
        size=10.5, color=MUTED)

    age_rows = [
        ('Chicken pox',  '2 – 12 yrs',   TEAL),
        ('Acne',         '13 – 25 yrs',   GREEN),
        ('Diabetes',     '35 – 75 yrs',   AMBER),
        ('Heart attack', '45 – 80 yrs',   RED),
        ('Malaria',      '5  – 65 yrs',   BLUE),
    ]
    for i, (d, r, c) in enumerate(age_rows):
        y = MT + 2.48 + i * 0.48
        box(slide, ML + 6.34, y, 5.65, 0.38,
            fill=RGBColor(0x12, 0x1C, 0x2A), stroke=None, rounded=True)
        txt(slide, d, ML + 6.52, y + 0.07, 3.2, 0.24,
            size=10.5, color=SOFT)
        txt(slide, r, ML + 9.9, y + 0.07, 1.8, 0.24,
            size=10.5, color=c, bold=True, align=PP_ALIGN.RIGHT)
    anim.add(rc, 280, 500, 'fade')

    anim.apply()
    transition(slide, 'wipe', 650)


# ── 5. PIPELINE ───────────────────────────────────────────────────────────────
def s_pipeline():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, BLUE)
    header_block(slide, 'Methodology', 'End-to-End Pipeline',
                 'A linear flow from raw CSVs to model-ready feature matrix and validated predictions.')
    footer(slide)
    anim = Animator(slide)

    steps = [
        # Row 1
        ('01', 'Load CSVs',           'disease_raw.csv\n+ symptom_severity.csv',  ORANGE),
        ('02', 'Build Feature Matrix','136-col severity-\nweighted symptom cols',  TEAL),
        ('03', 'Add Age Feature',     'Sample from epidemio-\nlogical age range',  BLUE),
        # Row 2
        ('04', '50% Symptom Dropout', 'Simulate incomplete\nclinical reporting',   AMBER),
        ('05', 'Label Encode',        'prognosis → integer\nclass index (41 cls)', PURPLE),
        ('06', 'Stratified Split',    '80/20 · 3,940 train\n985 test samples',     GREEN),
        # Row 3
        ('07', 'Variance Filter',     'Remove zero-var cols\n→ 135 final features',RED),
        ('08', 'Train 4 Classifiers', 'DT · RF(200) · XGB\n· Stacking Ensemble',  ORANGE),
        ('09', 'Evaluate',            'Accuracy · F1 · Precision\nRecall · CM ✓',  LIME),
    ]

    sw = (CW - 0.60) / 3   # step card width ≈ 3.91"
    sh = 1.52               # reduced: row2_bot=MT+1.78+1.52=4.52; row3_bot=MT+3.50+1.52=6.24
    row_y = [MT + 0.05, MT + 1.78, MT + 3.50]  # row_gap = 0.26" each

    for i, (num, title, sub, accent) in enumerate(steps):
        row, col = divmod(i, 3)
        x = ML + col * (sw + 0.30)
        y = row_y[row]
        c = card(slide, x, y, sw, sh, accent=accent, fill=PANEL, shadow=True)

        oval(slide, x + 0.20, y + 0.22, 0.40, 0.40, fill=accent)
        txt(slide, num, x + 0.20, y + 0.28, 0.40, 0.28,
            size=9, color=BG, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, title,  x + 0.78, y + 0.18, sw - 0.95, 0.30,
            size=12, bold=True, color=WHITE, font='Segoe UI Semibold')
        txt_multi(slide, sub.split('\n'), x + 0.22, y + 0.60, sw - 0.35, 0.82,
                  size=9.5, color=MUTED, para_space=2)

        # arrow connector (right side of non-last-in-row)
        if col < 2:
            connector(slide, x + sw + 0.04, y + sh / 2,
                      x + sw + 0.26, y + sh / 2, ORANGE, 1.5)

        anim.add(c, i * 100, 430, 'fade')

    # bottom note — row3_bot=6.24, note_top=6.46 → gap=0.22" ✓
    note = box(slide, ML, SH - MB - 0.62, CW, 0.36,
               fill=PANEL2, stroke=BORDER, rounded=True)
    txt(slide,
        'Design goal: every preprocessing step preserves interpretable medical signal.',
        ML + 0.22, SH - MB - 0.58, CW - 0.44, 0.28,
        size=10, color=MUTED, align=PP_ALIGN.CENTER)

    anim.apply()
    transition(slide, 'push', 650)


# ── 6. FOUR MODELS ────────────────────────────────────────────────────────────
def s_models():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, PURPLE)
    header_block(slide, 'Models', 'Four Classifier Comparison',
                 'A progression from interpretable baseline to meta-learning stacked ensemble.')
    footer(slide)
    anim = Animator(slide)

    mw = (CW - 0.75) / 4   # ≈ 2.895" per card
    models = [
        ('Decision Tree',       '85.58%', 'Baseline',
         ['criterion = gini', 'No depth limit', 'Fast & interpretable',
          'Overfits on dropout noise'],
         RED),
        ('Random Forest\n(200 trees)', '93.30%', 'Stable',
         ['200 estimators', 'Parallel bagging', 'Feature importance',
          'Smooths dropout noise'],
         TEAL),
        ('XGBoost',             '91.47%', 'Boosted',
         ['100 trees, lr=0.1', 'max_depth=6', 'mlogloss eval metric',
          'Sequential residuals'],
         BLUE),
        ('Stacking\nEnsemble',  '93.50%', '★ BEST',
         ['DT + RF + XGB base', 'LR meta-learner', 'passthrough=True',
          '5-fold stratified CV'],
         ORANGE),
    ]

    for i, (name, acc, tag, bullets, accent) in enumerate(models):
        x = ML + i * (mw + 0.25)
        ch = 5.38 if tag == '★ BEST' else 5.08
        c = card(slide, x, MT, mw, ch, accent=accent, fill=PANEL, shadow=True)

        # accent header
        box(slide, x, MT, mw, 0.85, fill=accent, rounded=True, stroke=None)
        # split name if \n
        name_lines = name.split('\n')
        txt(slide, name_lines[0], x, MT + 0.02, mw, 0.38,
            size=13, bold=True, color=BG,
            align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
        if len(name_lines) > 1:
            txt(slide, name_lines[1], x, MT + 0.38, mw, 0.30,
                size=10, bold=True, color=BG,
                align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
        # accuracy
        txt(slide, acc, x, MT + 0.96, mw, 0.50,
            size=22, color=accent, bold=True,
            align=PP_ALIGN.CENTER, font='Segoe UI Black')
        # tag pill
        tag_pill(slide, tag, x + mw / 2 - 0.50, MT + 1.48, accent)
        # bullets
        for j, b in enumerate(bullets):
            by = MT + 2.08 + j * 0.60
            box(slide, x + 0.18, by, mw - 0.28, 0.50,
                fill=RGBColor(0x16, 0x1E, 0x30), stroke=None, rounded=True)
            txt(slide, f'▸  {b}', x + 0.28, by + 0.10, mw - 0.40, 0.30,
                size=10, color=SOFT)
        # "RECOMMENDED" badge for stacking
        if tag == '★ BEST':
            rb = box(slide, x + 0.20, MT + ch - 0.55, mw - 0.40, 0.34,
                     fill=ORANGE, rounded=True, stroke=None, shadow=True)
            txt(slide, 'RECOMMENDED', x + 0.20, MT + ch - 0.50, mw - 0.40, 0.24,
                size=9, color=BG, bold=True, align=PP_ALIGN.CENTER)

        anim.add(c, i * 140, 520, 'fly')

    anim.apply()
    transition(slide, 'fade', 700)


# ── 7. RESULTS DASHBOARD ──────────────────────────────────────────────────────
def s_results():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, ORANGE)
    header_block(slide, 'Results', 'Performance Dashboard',
                 'Stacking Ensemble leads with 93.50% accuracy — 7.92 points above the Decision Tree baseline.')
    footer(slide)
    anim = Animator(slide)

    # ── Left: bar chart panel ─────────────────────────────────────────────────
    lp = card(slide, ML, MT, 7.30, 5.20, accent=ORANGE, fill=PANEL, shadow=True)
    txt(slide, 'Accuracy Ranking', ML + 0.22, MT + 0.18, 4.5, 0.34,
        size=16, bold=True, color=WHITE, font='Segoe UI Semibold')

    bars = [('Decision Tree',        85.58, RED),
            ('XGBoost',              91.47, BLUE),
            ('Random Forest (200)',  93.30, GREEN),
            ('Stacking Ensemble',    93.50, ORANGE)]
    for i, (lbl, pct, c) in enumerate(bars):
        # max_bar=3.20: x(0.72)+lw(2.6)+gap(0.12)+bar(3.20)+gap(0.22)+val(0.85)=7.71 < panel_right(7.80)
        progress_row(slide, lbl, pct, ML + 0.22, MT + 0.72 + i * 0.96,
                     6.68, c, max_bar=3.20)

    # delta KPIs
    metric_tile(slide, ML + 0.22, MT + 4.45, 2.30, 0.68,
                '+7.92', 'accuracy pts vs DT baseline', AMBER)
    metric_tile(slide, ML + 2.72, MT + 4.45, 2.30, 0.68,
                '4', 'models benchmarked', TEAL)
    metric_tile(slide, ML + 5.22, MT + 4.45, 1.88, 0.68,
                '985', 'test samples', BLUE)

    anim.add(lp, 0, 550, 'fade')

    # ── Right: metrics table ──────────────────────────────────────────────────
    rp = card(slide, ML + 7.64, MT, 4.60, 5.20, accent=AMBER, fill=PANEL, shadow=True)

    # table header
    box(slide, ML + 7.64, MT, 4.60, 0.48,
        fill=RGBColor(0x28, 0x22, 0x0A), stroke=None, rounded=True)
    headers = ['Model', 'Acc.', 'Prec.', 'Rec.', 'F1']
    col_x   = [ML + 7.84, ML + 9.38, ML + 10.22, ML + 10.90, ML + 11.58]
    col_w   = [1.40, 0.70, 0.65, 0.65, 0.60]
    for j, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        txt(slide, h, cx, MT + 0.09, cw, 0.28,
            size=9.5, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

    rows = [('Decision Tree',       '85.58', '0.86', '0.86', '0.86', RED),
            ('XGBoost',             '91.47', '0.92', '0.91', '0.91', BLUE),
            ('Random Forest',       '93.30', '0.94', '0.93', '0.93', GREEN),
            ('Stacking (★ Best)', '93.50', '0.94', '0.93', '0.93', ORANGE)]
    for ri, (*vals, c) in enumerate(rows):
        ry = MT + 0.60 + ri * 0.82
        bg = RGBColor(0x20, 0x16, 0x06) if ri == 3 else RGBColor(0x14, 0x19, 0x26)
        box(slide, ML + 7.84, ry, 4.18, 0.68,
            fill=bg, stroke=None, rounded=True)
        for j, (v, cx, cw) in enumerate(zip(vals, col_x, col_w)):
            txt(slide, v, cx, ry + 0.16, cw, 0.36,
                size=10 if ri < 3 else 10.5,
                color=c if ri == 3 or j == 0 else SOFT,
                bold=(ri == 3),
                align=PP_ALIGN.CENTER)
        if ri < 3:
            connector(slide, ML + 7.84, ry + 0.70,
                      ML + 12.02, ry + 0.70, BORDER, 0.5)

    # highlight strip
    box(slide, ML + 7.64, MT + 4.48, 4.60, 0.56,
        fill=RGBColor(0x26, 0x18, 0x06), stroke=ORANGE, sw=1.2, rounded=True)
    txt(slide,
        '★  Stacking Ensemble 93.50%  —  macro F1 0.93  —  BEST MODEL',
        ML + 7.84, MT + 4.58, 4.18, 0.36,
        size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    anim.add(rp, 200, 500, 'fade')
    anim.apply()
    transition(slide, 'fade', 700)


# ── 8. STACKING DEEP DIVE ─────────────────────────────────────────────────────
def s_stacking():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, AMBER)
    header_block(slide, 'Architecture', 'Stacking Ensemble — Deep Dive',
                 'Base learners produce probability outputs; the meta-learner blends them + raw features.')
    footer(slide)
    anim = Animator(slide)

    # ── Architecture diagram (fits in MT → MT+2.80) ──────────────────────────
    # Bases spaced at 0.82" intervals with h=0.65; base3 bottom = MT+0.20+2*0.80+0.65 = MT+2.45
    # Input card centre-y aligns with middle base learner centre
    bases = [('Decision\nTree', RED,  MT + 0.20),
             ('Random\nForest', TEAL, MT + 1.00),
             ('XGBoost',        BLUE, MT + 1.80)]

    # Input card: vertically centred to middle base
    mid_base_cy = MT + 1.00 + 0.325   # = MT+1.325
    ic_h = 0.82
    ic_y = mid_base_cy - ic_h / 2     # = MT+0.915
    ic = card(slide, ML, ic_y, 1.90, ic_h, accent=BLUE, fill=PANEL2)
    txt(slide, 'Input', ML + 0.10, ic_y + 0.10, 1.70, 0.30,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, '135 features', ML + 0.10, ic_y + 0.44, 1.70, 0.24,
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

    for name, c, y in bases:
        bc = card(slide, ML + 2.45, y, 2.50, 0.65, accent=c, fill=PANEL2)
        lines = name.split('\n')
        y_off = 0.04 if len(lines) > 1 else 0.18
        txt(slide, lines[0], ML + 2.45, y + y_off, 2.50, 0.26,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
        if len(lines) > 1:
            txt(slide, lines[1], ML + 2.45, y + 0.32, 2.50, 0.24,
                size=10, color=SOFT, align=PP_ALIGN.CENTER)
        # connector: from input right-centre to base left-centre
        connector(slide, ML + 1.93, mid_base_cy, ML + 2.42, y + 0.325, c, 1.1)
        anim.add(bc, 0, 450, 'fade')

    # base3 bottom = MT + 1.80 + 0.65 = MT + 2.45
    # Passthrough arrow (input → meta, above the bases)
    meta_cy = bases[1][2] + 0.325   # same as middle base centre
    connector(slide, ML + 1.93, mid_base_cy, ML + 5.65, meta_cy, MUTED, 1.0)
    txt(slide, 'passthrough=True', ML + 3.0, mid_base_cy - 0.24, 2.5, 0.22,
        size=8, color=MUTED, italic=True)

    # Connectors base right → meta left
    meta_y = bases[1][2]   # MT+1.00 — meta card top aligns with middle base
    for _, c, y in bases:
        connector(slide, ML + 4.98, y + 0.325, ML + 5.62, meta_cy, c, 1.1)

    # Meta-learner
    mc = card(slide, ML + 5.65, meta_y - 0.10, 3.0, 1.00, accent=AMBER, fill=PANEL)
    txt(slide, 'Meta-Learner', ML + 5.65, meta_y - 0.04, 3.0, 0.28,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Segoe UI Semibold')
    txt(slide, 'Logistic Regression', ML + 5.65, meta_y + 0.26, 3.0, 0.22,
        size=10, color=AMBER, align=PP_ALIGN.CENTER)
    txt(slide, 'C=0.1  |  saga  |  max_iter=3000', ML + 5.65, meta_y + 0.50, 3.0, 0.22,
        size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    anim.add(mc, 200, 500, 'fade')

    connector(slide, ML + 8.68, meta_cy, ML + 9.30, meta_cy, LIME, 1.5)

    # Output
    out_y = meta_y - 0.10
    oc = card(slide, ML + 9.32, out_y, 2.45, 1.00, accent=LIME, fill=PANEL2)
    txt(slide, 'Prediction', ML + 9.32, out_y + 0.10, 2.45, 0.28,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, 'Final disease label', ML + 9.32, out_y + 0.44, 2.45, 0.24,
        size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Key-points row ───────────────────────────────────────────────────────
    # base3 bottom = MT+2.45=3.67  →  kps start at MT+2.65=3.87 (0.20" safety gap)
    kps = [
        ('5-Fold CV',
         'Out-of-fold predictions during training prevent data leakage to the meta-learner.',
         AMBER),
        ('passthrough=True',
         'Meta-learner sees raw features AND base model outputs — corrects cases all 3 models miss.',
         ORANGE),
        ('predict_proba',
         'Passes full probability distributions, not just hard labels — richer meta-input.',
         TEAL),
        ('+0.20% over RF',
         'Meta-learner recovers residual errors on hard disease-boundary samples.',
         GREEN),
    ]
    kpw = (CW - 0.75) / 4   # 2.895"
    kp_y = MT + 2.65         # kps top = 3.87;  bottom = 3.87+1.72=5.59 < footer(7.08) ✓
    kp_h = 1.72
    for i, (kw, desc, c) in enumerate(kps):
        x = ML + i * (kpw + 0.25)
        kc = card(slide, x, kp_y, kpw, kp_h, accent=c, fill=PANEL, shadow=True)
        txt(slide, kw,   x + 0.20, kp_y + 0.12, kpw - 0.30, 0.30,
            size=11.5, bold=True, color=c, font='Segoe UI Semibold')
        txt(slide, desc, x + 0.20, kp_y + 0.52, kpw - 0.30, 1.10,
            size=10, color=MUTED)
        anim.add(kc, 400 + i * 120, 450, 'fade')

    anim.apply()
    transition(slide, 'push', 650)


# ── 9. FEATURE IMPORTANCE ─────────────────────────────────────────────────────
def s_features():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, LIME)
    header_block(slide, 'Explainability', 'Feature Importance  (Random Forest)',
                 'Top symptoms driving disease discrimination — and confirmation that age adds real predictive signal.')
    footer(slide)
    anim = Animator(slide)

    feats = [
        ('yellowing_of_eyes',    0.048, ORANGE),
        ('dark_urine',           0.044, ORANGE),
        ('sweating',             0.041, TEAL),
        ('chills',               0.038, TEAL),
        ('age',                  0.035, RED),     # highlighted
        ('loss_of_appetite',     0.033, TEAL),
        ('fatigue',              0.031, TEAL),
        ('high_fever',           0.029, TEAL),
        ('vomiting',             0.027, TEAL),
        ('nausea',               0.025, TEAL),
    ]
    max_v = 0.050
    bw    = 4.80   # max bar width
    lw    = 2.65   # label col width
    vw    = 0.70   # value col width

    # chart panel
    cp = card(slide, ML, MT, 7.60, 5.50, accent=ORANGE, fill=PANEL, shadow=True)
    txt(slide, 'Top 10 Most Discriminative Features', ML + 0.22, MT + 0.18, 7.0, 0.34,
        size=15, bold=True, color=WHITE, font='Segoe UI Semibold')

    for i, (name, val, c) in enumerate(feats):
        y = MT + 0.66 + i * 0.46
        label_disp = name.replace('_', ' ').title()
        is_age = (name == 'age')
        txt(slide, label_disp,
            ML + 0.22, y + 0.03, lw, 0.32,
            size=10, color=ORANGE if is_age else SOFT,
            bold=is_age)
        # track
        box(slide, ML + lw + 0.35, y + 0.06, bw, 0.14,
            fill=RGBColor(0x1A, 0x22, 0x34), rounded=True, stroke=None)
        # fill
        box(slide, ML + lw + 0.35, y + 0.06, bw * val / max_v, 0.14,
            fill=c, rounded=True, stroke=None)
        # value
        txt(slide, f'{val:.3f}',
            ML + lw + bw + 0.44, y + 0.02, vw, 0.28,
            size=9, color=c, bold=True)

    anim.add(cp, 0, 550, 'fade')

    # Insights panel
    ip = card(slide, ML + 7.92, MT, 4.28, 5.50, accent=AMBER, fill=PANEL, shadow=True)
    txt(slide, 'Key Insights', ML + 8.12, MT + 0.18, 3.85, 0.34,
        size=15, bold=True, color=WHITE, font='Segoe UI Semibold')

    insights = [
        (ORANGE, 'Jaundice cluster',
         'Yellowing of eyes + dark urine rank #1 and #2 — strong hepatitis discriminators.'),
        (RED,    'Age ranks #5',
         'Confirms epidemiological prior knowledge adds real signal beyond symptoms alone.'),
        (TEAL,   'Fever-adjacent features',
         'Sweating, chills, high fever rank high but are less unique — shared across many diseases.'),
        (BLUE,   'Severity matters',
         'Continuous severity weights give each symptom a gradient signal, not just binary presence.'),
        (GREEN,  'Clinical alignment',
         'The importance ranking matches known clinical diagnostic hierarchies for these diseases.'),
    ]
    for i, (c, heading, body) in enumerate(insights):
        y = MT + 0.72 + i * 0.94
        oval(slide, ML + 8.12, y + 0.06, 0.18, 0.18, fill=c)
        txt(slide, heading, ML + 8.46, y + 0.04, 3.52, 0.26,
            size=10.5, bold=True, color=c, font='Segoe UI Semibold')
        txt(slide, body,    ML + 8.46, y + 0.30, 3.52, 0.55,
            size=9, color=MUTED)
        if i < 4:
            connector(slide, ML + 8.12, y + 0.88,
                      ML + 11.92, y + 0.88, BORDER, 0.4)

    anim.add(ip, 200, 500, 'fade')
    anim.apply()
    transition(slide, 'wipe', 650)


# ── 10. DISCUSSION ────────────────────────────────────────────────────────────
def s_discussion():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, GREEN)
    header_block(slide, 'Discussion', 'Strengths & Limitations',
                 'What the numbers mean and where the current system can still improve.')
    footer(slide)
    anim = Animator(slide)

    pw = (CW - 0.40) / 2

    # Strengths
    sc = card(slide, ML, MT, pw, 5.38, accent=GREEN, fill=PANEL, shadow=True)
    box(slide, ML, MT, pw, 0.50, fill=GREEN, rounded=True, stroke=None)
    txt(slide, '✓  What Works Well', ML + 0.22, MT + 0.08, pw - 0.30, 0.34,
        size=14, bold=True, color=BG, font='Segoe UI Semibold')

    strengths = [
        ('Severity encoding',      'Outperforms binary encoding — each symptom contributes a continuous clinical weight.'),
        ('Age feature (#5 rank)',   'Epidemiological priors add real discriminative power, confirmed by importance scores.'),
        ('50% dropout training',   'Models are robust to missing symptoms at test time — realistic clinical behaviour.'),
        ('Balanced dataset',       'All 41 classes have ~120 records → accuracy is a meaningful metric here.'),
        ('Stacking passthrough',   'Meta-learner corrects cases where all base models agree but are wrong.'),
    ]
    for i, (h, d) in enumerate(strengths):
        y = MT + 0.65 + i * 0.92
        box(slide, ML + 0.20, y, pw - 0.30, 0.82,
            fill=RGBColor(0x0C, 0x1A, 0x10), stroke=None, rounded=True)
        txt(slide, h, ML + 0.35, y + 0.08, pw - 0.52, 0.26,
            size=10.5, bold=True, color=GREEN, font='Segoe UI Semibold')
        txt(slide, d, ML + 0.35, y + 0.36, pw - 0.52, 0.40,
            size=9.5, color=MUTED)
    anim.add(sc, 0, 520, 'fade')

    # Limitations
    lc = card(slide, ML + pw + 0.40, MT, pw, 5.38, accent=RED, fill=PANEL, shadow=True)
    box(slide, ML + pw + 0.40, MT, pw, 0.50, fill=RED, rounded=True, stroke=None)
    txt(slide, '✗  Current Limitations', ML + pw + 0.62, MT + 0.08, pw - 0.30, 0.34,
        size=14, bold=True, color=BG, font='Segoe UI Semibold')

    limits = [
        ('Structured checklist',   'Patients describe symptoms in natural language — not ticked checkboxes.'),
        ('Hepatitis variants',      'B / C / D / E still confused — lab tests needed for final differentiation.'),
        ('No hyperparameter tuning','Optuna / GridSearchCV could push accuracy beyond current levels.'),
        ('Synthetic-style data',    'Clean, structured labels — real clinical records are significantly noisier.'),
        ('No vital signs',          'No gender, blood pressure, temperature, or lab values as features yet.'),
    ]
    for i, (h, d) in enumerate(limits):
        y = MT + 0.65 + i * 0.92
        box(slide, ML + pw + 0.60, y, pw - 0.30, 0.82,
            fill=RGBColor(0x1A, 0x0C, 0x0C), stroke=None, rounded=True)
        txt(slide, h, ML + pw + 0.75, y + 0.08, pw - 0.52, 0.26,
            size=10.5, bold=True, color=RED, font='Segoe UI Semibold')
        txt(slide, d, ML + pw + 0.75, y + 0.36, pw - 0.52, 0.40,
            size=9.5, color=MUTED)
    anim.add(lc, 160, 520, 'fade')

    anim.apply()
    transition(slide, 'fade', 700)


# ── 11. FUTURE SCOPE ──────────────────────────────────────────────────────────
def s_future():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    top_strip(slide, BLUE)
    header_block(slide, 'Roadmap', 'Future Scope & Deployment Path',
                 'From academic prototype to a deployable clinical decision-support product.')
    footer(slide)
    anim = Animator(slide)

    future = [
        ('Hyperparameter Tuning',  'Optuna / GridSearchCV on XGB depth and LR regularisation constant C.',            BLUE,   0),
        ('Neural Networks',        'MLP or Transformer encoder for higher-order symptom co-occurrence patterns.',      PURPLE, 1),
        ('SHAP Explainability',    'Per-patient "why" explanations — critical for clinical trust and auditability.',   TEAL,   2),
        ('Web Deployment',         'FastAPI backend + React symptom-input form → working prototype in a browser.',     GREEN,  3),
        ('Richer Patient Data',    'Add gender, blood pressure, temperature, and basic lab values (CBC, LFT).',        AMBER,  4),
        ('Clinical NLP',           'Extract symptoms from free-text discharge summaries via NLP / LLM pipeline.',      ORANGE, 5),
    ]

    fw = (CW - 0.50) / 3
    for i, (title, desc, c, idx) in enumerate(future):
        row, col = divmod(idx, 3)
        x = ML + col * (fw + 0.25)
        y = MT + row * 2.42

        fc = card(slide, x, y, fw, 2.12, accent=c, fill=PANEL, shadow=True)
        num_badge(slide, idx + 1, x + 0.22, y + 0.22, c)
        txt(slide, title, x + 0.82, y + 0.22, fw - 0.98, 0.32,
            size=13, bold=True, color=WHITE, font='Segoe UI Semibold')
        txt(slide, desc, x + 0.22, y + 0.70, fw - 0.36, 1.26,
            size=10.5, color=MUTED)
        anim.add(fc, i * 110, 450, 'fly')

    # Applications callout
    ac = box(slide, ML, SH - MB - 0.56, CW, 0.42,
             fill=RGBColor(0x10, 0x1A, 0x2A), stroke=BLUE, sw=1.2, rounded=True)
    txt(slide,
        'Applications:  GP clinics  ·  Remote health workers  ·  Symptom checker apps  '
        '·  Hospital triage  ·  Medical education tools',
        ML + 0.25, SH - MB - 0.50, CW - 0.50, 0.32,
        size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    anim.add(ac, 700, 450)

    anim.apply()
    transition(slide, 'push', 650)


# ── 12. CONCLUSION ────────────────────────────────────────────────────────────
def s_conclusion():
    slide = prs.slides.add_slide(BLANK)
    slide_bg(slide)
    anim = Animator(slide)

    # Full-width top gradient
    box(slide, 0, 0, SW, 0.055, fill=ORANGE, rounded=False)
    box(slide, 0, 0.055, SW, 1.10, fill=RGBColor(0x14, 0x0C, 0x04), rounded=False)

    tag_pill(slide, 'FINAL TAKEAWAY', ML, 0.18)
    txt(slide, 'From Symptoms to Confident Prediction', ML, 0.50, 9.0, 0.62,
        size=30, bold=True, color=WHITE, font='Segoe UI Semibold')

    # Accuracy recap row
    results = [('Decision Tree',       '85.58 %', RED),
               ('XGBoost',             '91.47 %', BLUE),
               ('Random Forest (200)', '93.30 %', GREEN),
               ('Stacking Ensemble',   '93.50 %', ORANGE)]
    rw = (CW - 0.75) / 4
    for i, (m, a, c) in enumerate(results):
        x = ML + i * (rw + 0.25)
        rc = box(slide, x, 1.30, rw, 0.98,
                 fill=RGBColor(0x14, 0x19, 0x26), stroke=c, sw=1.5, rounded=True)
        txt(slide, m, x, 1.35, rw, 0.34, size=9.5, color=SOFT, align=PP_ALIGN.CENTER)
        txt(slide, a, x, 1.67, rw, 0.46, size=19, color=c, bold=True,
            align=PP_ALIGN.CENTER, font='Segoe UI Black')
        anim.add(rc, i * 120, 480, 'fade')

    # Star badge on stacking
    box(slide, ML + 3 * (rw + 0.25), 1.26, rw, 0.06, fill=ORANGE, rounded=False)

    # Key takeaways
    txt(slide, 'Key Takeaways', ML, 2.50, 4.0, 0.36,
        size=15, bold=True, color=ORANGE)
    connector(slide, ML, 2.85, ML + 8.8, 2.85, BORDER, 0.8)

    takes = [
        'Classical ML with thoughtful feature engineering stays highly competitive for structured medical data.',
        'Severity-weighted symptoms + epidemiological age outperform plain binary symptom encoding.',
        'Stacking Ensemble with passthrough=True squeezes residual errors that no single model corrects alone.',
        '50% dropout simulation makes the system robust to missing symptoms in real-world usage.',
        'Next step: web deployment with SHAP explanations for clinical interpretability.',
    ]
    for i, t in enumerate(takes):
        y = 3.00 + i * 0.74
        tb = box(slide, ML, y, 8.80, 0.62,
                 fill=RGBColor(0x12, 0x16, 0x24), stroke=BORDER, rounded=True)
        oval(slide, ML + 0.16, y + 0.20, 0.22, 0.22, fill=ORANGE)
        txt(slide, t, ML + 0.52, y + 0.12, 8.12, 0.40,
            size=11, color=SOFT)
        anim.add(tb, 200 + i * 150, 430, 'fade')

    # Right panel
    qp = card(slide, 9.95, 1.30, 2.90, 4.60, accent=ORANGE, fill=PANEL2, shadow=True)
    txt(slide, 'Questions', 9.95, 2.10, 2.90, 0.55,
        size=26, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font='Segoe UI Black')
    txt(slide, 'Thank you', 9.95, 2.72, 2.90, 0.36,
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    connector(slide, 10.30, 3.22, 12.52, 3.22, ORANGE, 1.4)
    txt(slide, 'Disease Prediction System', 9.95, 3.40, 2.90, 0.26,
        size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, 'MLC-2  |  CSE 3968  |  ITER', 9.95, 3.76, 2.90, 0.24,
        size=9, color=MUTED, align=PP_ALIGN.CENTER)
    txt(slide, 'SOA University  |  May 2026', 9.95, 4.00, 2.90, 0.24,
        size=9, color=MUTED, align=PP_ALIGN.CENTER)
    metric_tile(slide, 10.18, 4.48, 2.44, 0.96,
                '93.50%', 'Best model accuracy', ORANGE)
    anim.add(qp, 80, 520, 'fade')

    footer(slide)
    anim.apply()
    transition(slide, 'fade', 900)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD DECK
# ═══════════════════════════════════════════════════════════════════════════════
builders = [
    s_cover, s_agenda, s_problem, s_dataset,
    s_pipeline, s_models, s_results, s_stacking,
    s_features, s_discussion, s_future, s_conclusion,
]

for fn in builders:
    fn()

prs.save(str(OUT))
print(f'Saved → {OUT}')
print(f'Slides : {len(prs.slides)}')
print('Theme  : Premium dark  ·  Orange/Teal/Navy')
